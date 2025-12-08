from __future__ import annotations
import base64
import datetime
import io
import statistics
from typing import Optional, List, Dict, Any, Iterable, Union
import os, uuid
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, RGBColor
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from pptx.util import Inches, Pt
from supabase import Client, create_client

SUPABASE_URL = "http://103.175.218.139:8000"
SUPABASE_KEY = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJyb2xlIjoic2VydmljZV9yb2xlIiwiaXNzIjoic3VwYWJhc2UiLCJpYXQiOjE3NTc4NTU5MzIsImV4cCI6MjA3MzIxNTkzMn0.pAnC8zBC7ewuIkvBGG_ds1lxyXsAxsCA3OXelt0XXWI"

supabase : Client = create_client(SUPABASE_URL, SUPABASE_KEY)

def upload_file_and_get_url_public(file_bytes: bytes, file_name:str, email:str) -> Optional[str]:
    """
    Mengunggah file stream ke Supabase Storage dan mengembalikan URL publiknya.
    """
    # Inisialisasi Klien Supabase
    
    try:    
        # Unggah byte stream langsung
        print(f"Mengunggah '{file_name}' ke bucket 'testing'...")
        path = f"{email}/{file_name}"
        supabase.storage.from_("testing").upload(
            path=path,
            file=file_bytes, # Menggunakan byte stream langsung
            # file_options={"content-type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "upsert": "true"}
        )
        print("Unggah berhasil.")

        # Dapatkan URL publik
        print("Mengambil URL unduhan...")
        public_url = supabase.storage.from_("testing").get_public_url(path)

        print(f"URL Unduhan: {public_url}")
        
        return public_url

    except Exception as e:
        print(f"Terjadi kesalahan saat berinteraksi dengan Supabase: {e}")
        return None

def download_file_as_bytes(
    bucket_name: str,
    file_path: str,
) -> Optional[bytes]:
    """
    Mengunduh file dari Supabase Storage sebagai bytes.
    - bucket_name: nama bucket (mis. 'testing')
    - file_path: path lengkap objek (mis. 'folder/dokumen.docx')
    Return: bytes konten file atau None jika gagal.
    """
    try:
        data: bytes = supabase.storage.from_(bucket_name).download(file_path)
        print("\033[31mDOWNLOAD FILE\033[0m", data)
        return data
    except Exception as e:
        # Log sesuai kebutuhan
        print(f"Download gagal: {e}")
        return None

def generate_chart(data: dict, title: str) -> str:
    """
    Buat chart dari data, simpan jadi PNG base64.
    """
    import matplotlib.pyplot as plt
    import io, base64

    plt.figure()
    plt.plot(data["x"], data["y"])
    plt.title(title)
    plt.xlabel("Bulan")
    plt.ylabel("Penjualan")

    buffer = io.BytesIO()
    plt.savefig(buffer, format='png')
    buffer.seek(0)
    return base64.b64encode(buffer.getvalue()).decode('utf-8')


def create_docx(
    spec: Dict[str, Any],
    filename: Optional[str] = None,
    email: str = "default"

) -> Dict[str, Any]:
    """
    create_docx: Membuat berkas Word (.docx) dari spesifikasi terstruktur.

    Gunakan saat: agent sudah memiliki outline/konten final dan ingin menghasilkan file .docx.

    Input:

    spec (wajib, object):

        title (string), subtitle (string, opsional), metadata (object, opsional: author/subject/keywords/…), style (opsional),

        sections (wajib, array): tiap section boleh berisi heading, paragraphs[], bullets[], numbered[], table{headers[], rows[][]}, images[] (base64).

        filename (opsional, string): nama file; jika tanpa ekstensi akan ditambah .docx.

        Perilaku:

        Validasi: minimal 1 section; tabel harus kolom konsisten; gambar base64 divalidasi/ukuran dibatasi; metadata disetel ke core properties.

        Styling global ringan (line spacing/space after) tanpa merusak style daftar.

    Keluaran:

        object: { status: "ok"|"error", message: str,
        file: { name, path?, url?, bytes_len }, warnings: [..], issues: [{field, detail}] }
    """

    print("USER_ID_DOCS", email)
    result = {"status": "error", "message": "", "file": None, "warnings": [], "issues": []}
    try:
        # Validasi minimal
        title = spec.get("title") or "Dokumen Tanpa Judul"
        sections = spec.get("sections", [])
        if not isinstance(sections, list) or len(sections) == 0:
            result["issues"].append({"field": "sections", "detail": "Wajib diisi minimal 1"})
            result["message"] = "Input tidak valid"
            return result

        # Tentukan nama/path file deterministik
        base = filename or f"{uuid.uuid4()}.docx"   
        if not base.lower().endswith(".docx"):
            base += ".docx"
        print("FILE_NAME BUAT",uuid.uuid4())
        safe_base = f"{uuid.uuid4()}--{os.path.basename(base)}"

        doc = Document()

        # Core properties
        doc.core_properties.title = title
        metadata = spec.get("metadata") or {}
        if isinstance(metadata, dict):
            doc.core_properties.author = str(metadata.get("author", ""))[:255]
            doc.core_properties.subject = str(metadata.get("subject", ""))[:255]
            doc.core_properties.category = str(metadata.get("category", ""))[:255]
            doc.core_properties.keywords = str(metadata.get("keywords", ""))[:255]
            doc.core_properties.comments = str(metadata.get("comments", ""))[:255]

        # Judul dan subjudul
        title_para = doc.add_heading(title, level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        subtitle = spec.get("subtitle")
        if subtitle:
            sub_p = doc.add_paragraph(str(subtitle))
            sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if sub_p.runs:
                sub_p.runs[0].bold = True
            doc.add_paragraph()

        # Metadata visual opsional
        if metadata:
            meta_para = doc.add_paragraph()
            for k, v in metadata.items():
                run = meta_para.add_run(f"{k}: ")
                run.bold = True
                meta_para.add_run(f"{v}\n")
            doc.add_paragraph()

        # Sections
        for s in sections:
            heading = s.get("heading")
            if heading:
                doc.add_heading(str(heading), level=2)

            for para in s.get("paragraphs", []) or []:
                doc.add_paragraph(str(para))

            for b in s.get("bullets", []) or []:
                doc.add_paragraph(str(b), style="List Bullet")

            for n in s.get("numbered", []) or []:
                doc.add_paragraph(str(n), style="List Number")

            table_spec = s.get("table")
            if isinstance(table_spec, dict):
                headers = table_spec.get("headers") or []
                rows = table_spec.get("rows") or []
                if headers and all(isinstance(h, str) for h in headers):
                    table = doc.add_table(rows=1, cols=len(headers))
                    table.style = "Table Grid"
                    for i, h in enumerate(headers):
                        table.rows[0].cells[i].text = h
                    for row in rows:
                        row_cells = table.add_row().cells
                        for i, cell_text in enumerate(row[:len(headers)]):
                            row_cells[i].text = str(cell_text)
                    doc.add_paragraph()
                else:
                    result["warnings"].append("Table diabaikan: headers tidak valid")

            images = s.get("images") or []
            for img_b64 in images:
                try:
                    raw = base64.b64decode(img_b64, validate=True)
                    if len(raw) > 5 * 1024 * 1024:
                        result["warnings"].append("Gambar >5MB di-skip")
                        continue
                    stream = io.BytesIO(raw)
                    doc.add_picture(stream, width=Inches(5))
                    doc.add_paragraph()
                except Exception as e:
                    result["warnings"].append(f"Gagal memproses gambar: {e}")

        # Styling global ringan
        for p in doc.paragraphs:
            pf = p.paragraph_format
            pf.space_after = Pt(6)
            pf.line_spacing = 1.3

        # Simpan ke BytesIO dulu
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        data_bytes = buf.getvalue()

        # Jika punya uploader eksternal, panggil di sini
        url = None
        try:
            url = upload_file_and_get_url_public(data_bytes, safe_base, email)
        except Exception as e:
            result["warnings"].append(f"Upload gagal: {e}")

        result["status"] = "ok"
        result["message"] = "Dokumen berhasil dibuat"
        result["file"] = {"name": safe_base, "url": url, "bytes_len": len(data_bytes)}
        return result

    except Exception as e:
        result["status"] = "error"
        result["message"] = f"Gagal membuat dokumen: {e}"
        return result

def convert_docx_bytes_to_pdf_bytes(docx_bytes: bytes) -> bytes:
    """
    Konversi konten DOCX (bytes) ke PDF (bytes) menggunakan ReportLab.
    Menangani paragraf dan heading sederhana.
    """
    # Buka docx dari memori
    doc = Document(io.BytesIO(docx_bytes))

    # Siapkan buffer PDF in-memory
    pdf_buffer = io.BytesIO()
    doc_pdf = SimpleDocTemplate(pdf_buffer, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []

    # Render paragraf
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            story.append(Spacer(1, 12))
            continue
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            style = styles["Heading1"]
        else:
            # Heuristik: judul singkat kapital semua
            if len(text.split()) <= 5 and text.isupper():
                style = styles["Heading2"]
            else:
                style = styles["Normal"]
        story.append(Paragraph(text, style))
        story.append(Spacer(1, 8))

    # Bangun PDF ke buffer
    doc_pdf.build(story)
    pdf_buffer.seek(0)
    return pdf_buffer.read()

def convert_to_pdf_and_upload(
    source_bucket: str,
    source_path: str,          # contoh: "nama_user/filename.docx"
    target_bucket: str = "testing",
    target_path: Optional[str] = None,  # contoh: "nama_user/filename.pdf"
) -> Optional[str]:
    """
    Ambil DOCX dari Supabase Storage, konversi ke PDF (in-memory), 
    unggah ke Storage, lalu kembalikan public URL PDF.
    """
    # 1) Unduh DOCX sebagai bytes
    docx_bytes = download_file_as_bytes(bucket_name=source_bucket, file_path=source_path)
    if not docx_bytes:
        print(f"Gagal mengunduh sumber: {source_bucket}/{source_path}")
        return None

    # 2) Konversi DOCX -> PDF (in-memory)
    pdf_bytes = convert_docx_bytes_to_pdf_bytes(docx_bytes)

    # 3) Tentukan target path
    if not target_path:
        base, _ = os.path.splitext(source_path)
        target_path = f"{base}.pdf"

    # 4) Upload PDF ke Storage dengan content-type yang benar
    try:
        supabase.storage.from_(target_bucket).upload(
            path=target_path,
            file=pdf_bytes,
            file_options={
                "content-type": "application/pdf",
                "upsert": "true",  # set ke "false" jika ingin hindari overwrite
            },
        )
    except Exception as e:
        print(f"Upload PDF gagal: {e}")
        return None

    # 5) Ambil URL publik
    try:
        public_url = supabase.storage.from_(target_bucket).get_public_url(target_path)
        print("UPLOAD FILE CONVERT", public_url)
        return public_url
    except Exception as e:
        print(f"Gagal mengambil public URL: {e}")
        return None

def convert_to_pdf(filename: str,email: str = "default") -> str:
    """
    Konversi file DOCX di storage menjadi PDF dan kembalikan hasil terstruktur.

    Input:
    - filename: path/nama file DOCX pada storage (contoh: "user123/laporan.docx")
    - overwrite: izinkan menimpa file PDF yang sudah ada

    Output (dict):
    {
      "url": str|None 
    }
    """

    base_name = os.path.splitext(filename)[0]
    pdf_url = convert_to_pdf_and_upload(
    source_bucket="testing",
    source_path=f"{email}/{filename}",
    target_bucket="testing",
    target_path=f"{email}/{base_name}.pdf",
    )

    print("BBBBBBBBBERHASIL CONVERT", pdf_url)

    return {"url":pdf_url}

def create_pptx(
    spec: Dict[str, Any],
    email: str = "default",
    filename: Optional[str] = None,
) -> Dict[str, Any]:
    """
    create_pptx: Buat presentasi .pptx adaptif dari spesifikasi terstruktur.

    Gunakan saat: sudah ada tujuan/audiens/tema atau outline; jika slides kosong, struktur default dibuat berdasar purpose (laporan/proposal/umum).

    Input:

        spec (object, wajib)

            title (string), subtitle? (string)

            theme (string), purpose (string), audience (string)

            style_preset? ("formal"|"modern"|"creative")

            slides? (array): { title: string, content: string[], notes?: string[], layout?: "title+content"|"two-content"|"section" }

            num_slides? (int): pakai len(slides) atau 5 jika tidak ada

        filename? (string): tanpa ekstensi akan ditambah .pptx

    Perilaku:

    Auto‑generate slides bila kosong (berdasar purpose), isi placeholder singkat, batasi 3–6 bullet per slide, tambah cover otomatis, sisipkan notes jika ada, terapkan style_preset ringan.

    Keluaran:

    { status: "ok"|"error", message, file: { name, path?, url?, bytes_len }, meta: { slides_count, theme, purpose, audience, style_preset }, warnings: string[], issues: [{field, detail}] }

    """
    print("PPTX", filename)
    res = {"status": "error", "message": "", "file": None, "meta": {}, "warnings": [], "issues": []}
    try:
        title = (spec.get("title") or spec.get("theme") or "Presentasi").strip()
        subtitle = (spec.get("subtitle") or "Oleh: [Nama]").strip()
        theme = spec.get("theme") or "Topik Umum"
        purpose = (spec.get("purpose") or "").strip()
        audience = (spec.get("audience") or "").strip()
        style_preset = (spec.get("style_preset") or "formal").strip().lower()

        slides_data = spec.get("slides") or []
        num_slides = spec.get("num_slides") or (len(slides_data) if slides_data else 5)

        # Penentuan nama/path
        base = filename or f"{uuid.uuid4()}.pptx"
        if not base.lower().endswith(".pptx"):
            base += ".pptx"
        safe_base = os.path.basename(base)

        prs = Presentation()

        # Cover
        cover = prs.slides.add_slide(prs.slide_layouts[0])
        cover.shapes.title.text = title
        if len(cover.placeholders) > 1:
            cover.placeholders[1].text = subtitle

        # Generate default slides jika tidak ada
        def default_sections():
            if "laporan" in purpose.lower():
                return ["Ringkasan Eksekutif", "Kinerja", "Analisis Data", "Rencana", "Penutup"]
            if "proposal" in purpose.lower():
                return ["Latar Belakang", "Tujuan", "Metodologi", "Anggaran", "Penutup"]
            return ["Pendahuluan", "Analisis", "Strategi", "Implementasi", "Kesimpulan"]

        if not slides_data:
            sections = default_sections()[:max(1, num_slides)]
            slides_data = [
                {
                    "title": sec,
                    "content": [
                        f"Bahas {sec.lower()} terkait tema {theme.lower()}.",
                        f"Poin kunci untuk audiens {audience or 'umum'}.",
                        f"Implikasi dan rekomendasi singkat."
                    ]
                } for sec in sections
            ]

        # Peta layout sederhana
        layout_map = {
            "title+content": 1,
            "two-content": 3,
            "section": 2
        }

        # Tambah slide konten
        count = 0
        for s in slides_data[:num_slides]:
            layout_name = (s.get("layout") or "title+content").lower()
            layout_idx = layout_map.get(layout_name, 1)
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

            # Judul
            slide.shapes.title.text = s.get("title") or "Tanpa Judul"

            # Konten bullet
            if layout_idx in (1, 3):
                placeholder_idx = 1 if layout_idx == 1 else 1  # sederhana
                tf = slide.placeholders[placeholder_idx].text_frame
                tf.clear()
                bullets = [str(x) for x in (s.get("content") or []) if str(x).strip()]
                if not bullets:
                    bullets = ["[Isi poin utama]", "[Isi detail singkat]"]
                bullets = bullets[:6]
                for i, item in enumerate(bullets):
                    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                    p.text = item
                    p.level = 0 if i == 0 else 1

            # Notes
            notes = s.get("notes") or []
            if notes:
                ns = slide.notes_slide.notes_text_frame
                ns.text = "\n".join([str(n) for n in notes])

            count += 1

        # Style preset ringan (contoh: hanya cover title size)
        try:
            if style_preset == "modern":
                cover.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
            elif style_preset == "creative":
                cover.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)
            else:  # formal
                cover.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
        except Exception as e:
            res["warnings"].append(f"Gagal set font cover: {e}")

        # Simpan ke bytes
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        data_bytes = buf.getvalue()

        url = None
        try:
            url = upload_file_and_get_url_public(data_bytes, safe_base,email)
        except Exception as e:
            res["warnings"].append(f"Upload gagal: {e}")

        res["status"] = "ok"
        res["message"] = "PPTX berhasil dibuat"
        res["file"] = {"name": safe_base, "url": url, "bytes_len": len(data_bytes)}
        res["meta"] = {"slides_count": count, "theme": theme, "purpose": purpose, "audience": audience, "style_preset": style_preset}
        return res

    except Exception as e:
        res["status"] = "error"
        res["message"] = f"Gagal membuat PPTX: {e}"
        return res

